# -*- coding: utf-8 -*-
"""27 Markets — CRM back-office overview as a branded 16:9 PowerPoint deck.
Dark (black/red) theme matching the product. Diagrams rendered transparent."""
import os
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
from matplotlib.patches import FancyBboxPatch, FancyArrowPatch, Polygon
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

BASE = r'c:/Users/gowth/OneDrive/Desktop/27 trading'
DIAG = os.path.join(BASE, 'assets', 'diagrams')
os.makedirs(DIAG, exist_ok=True)
LOGO = r'C:/Users/gowth/Downloads/corecare consultancy.png'

# ───────────────────────── dark diagram engine ─────────────────────────
RED   = '#e11d2e'
PROC  = '#161a21'   # process box fill (visible on dark)
PEC   = '#39404c'   # process box edge
TXT   = '#ffffff'

def _anchor(cx, cy, w, h, tx, ty):
    dx, dy = tx - cx, ty - cy
    if abs(dy) >= abs(dx):
        return (cx, cy - h/2) if dy < 0 else (cx, cy + h/2)
    return (cx - w/2, cy) if dx < 0 else (cx + w/2, cy)

def draw_flow(nodes, edges, path, figsize):
    fig, ax = plt.subplots(figsize=figsize, dpi=200)
    fig.patch.set_alpha(0.0); ax.set_facecolor('none')
    for e in edges:
        a, b = e[0], e[1]
        label = e[2] if len(e) > 2 else None
        ax_, ay_, aw, ah = nodes[a][0], nodes[a][1], nodes[a][2], nodes[a][3]
        bx_, by_, bw, bh = nodes[b][0], nodes[b][1], nodes[b][2], nodes[b][3]
        p1 = _anchor(ax_, ay_, aw, ah, bx_, by_)
        p2 = _anchor(bx_, by_, bw, bh, ax_, ay_)
        ax.add_patch(FancyArrowPatch(p1, p2, arrowstyle='-|>', mutation_scale=15,
                     lw=1.8, color=RED, shrinkA=0, shrinkB=0, zorder=1))
        if label:
            mx, my = (p1[0]+p2[0])/2, (p1[1]+p2[1])/2
            ax.text(mx, my, label, fontsize=8, color='white', ha='center', va='center',
                    bbox=dict(boxstyle='round,pad=0.22', fc=PROC, ec=PEC, lw=0.8), zorder=3)
    for name, (cx, cy, w, h, label, kind) in nodes.items():
        if kind == 'decision':
            pts = [(cx, cy+h/2), (cx+w/2, cy), (cx, cy-h/2), (cx-w/2, cy)]
            ax.add_patch(Polygon(pts, closed=True, fc=PROC, ec=RED, lw=1.9, zorder=2))
            ax.text(cx, cy, label, ha='center', va='center', fontsize=8, color='white', zorder=3)
        elif kind == 'terminal':
            ax.add_patch(FancyBboxPatch((cx-w/2, cy-h/2), w, h,
                         boxstyle='round,pad=0.02,rounding_size=0.18', fc=RED, ec=RED, lw=1.5, zorder=2))
            ax.text(cx, cy, label, ha='center', va='center', fontsize=8.6, color='white', fontweight='bold', zorder=3)
        elif kind == 'group':
            ax.add_patch(FancyBboxPatch((cx-w/2, cy-h/2), w, h,
                         boxstyle='round,pad=0.02,rounding_size=0.10', fc=PROC, ec=RED, lw=1.5, ls='-', zorder=2))
            ax.text(cx, cy+h/2-0.18, label, ha='center', va='top', fontsize=8.8, color=RED, fontweight='bold', zorder=3)
        else:  # process
            ax.add_patch(FancyBboxPatch((cx-w/2, cy-h/2), w, h,
                         boxstyle='round,pad=0.02,rounding_size=0.12', fc=PROC, ec=PEC, lw=1.3, zorder=2))
            ax.text(cx, cy, label, ha='center', va='center', fontsize=8.4, color='white', zorder=3)
    xs=[v[0] for v in nodes.values()]; ys=[v[1] for v in nodes.values()]
    ws=[v[2] for v in nodes.values()]; hs=[v[3] for v in nodes.values()]
    ax.set_xlim(min(x-w/2 for x,w in zip(xs,ws))-0.4, max(x+w/2 for x,w in zip(xs,ws))+0.4)
    ax.set_ylim(min(y-h/2 for y,h in zip(ys,hs))-0.4, max(y+h/2 for y,h in zip(ys,hs))+0.4)
    ax.set_aspect('equal'); ax.axis('off')
    fig.savefig(path, bbox_inches='tight', transparent=True, dpi=200)
    plt.close(fig); return path

def P(cx,cy,label,w=2.6,h=0.9,kind='process'): return (cx,cy,w,h,label,kind)

# modules map
n = {
 'dash':(-3.5,1.7,3.0,1.2,'DASHBOARD\nKPIs · Activity\nAlerts','group'),
 'cli':(0,1.7,3.0,1.2,'CLIENTS\n360 View · Notes\nAccounts · History','group'),
 'lead':(3.5,1.7,3.0,1.2,'LEADS\nPipeline · Assign\nFollow-ups','group'),
 'kyc':(-3.5,-0.7,3.0,1.2,'KYC / COMPLIANCE\nReview Queue\nApprove · Reject','group'),
 'fin':(0,-0.7,3.0,1.2,'FINANCE\nDeposits · Withdrawal\nApprovals','group'),
 'sup':(3.5,-0.7,3.0,1.2,'SUPPORT\nTicket Inbox\nAssign · Reply','group'),
 'ib':(-3.5,-3.1,3.0,1.2,'PARTNERS / IB\nReferrals\nRebates','group'),
 'rep':(0,-3.1,3.0,1.2,'REPORTS\nRevenue · Funnel\nPerformance','group'),
 'staff':(3.5,-3.1,3.0,1.2,'STAFF & SETTINGS\nRoles · Audit Log\nAccess Control','group'),
}
draw_flow(n, [('dash','cli'),('cli','lead'),('cli','kyc'),('cli','fin'),('cli','sup'),
              ('cli','ib'),('dash','rep'),('staff','fin')], os.path.join(DIAG,'dk_modules.png'), (8.6,6.0))

GX,MW,MH,BRY = 3.0,2.6,1.15,-2.35
def hnode(i,label,kind='process',row=0): return P(i*GX,0 if row==0 else BRY,label,w=MW,h=MH,kind=kind)
def hdiamond(i,label,row=0): return (i*GX,0 if row==0 else BRY,2.7,1.35,label,'decision')

labels=['New\nLead','Assigned\nto Agent','Contacted\n/ Qualified','Registered\n+ Verified','Funded','Active\nClient']
n={f'c{i}':hnode(i,l,'terminal' if i in (0,5) else 'process') for i,l in enumerate(labels)}
draw_flow(n, [(f'c{i}',f'c{i+1}') for i in range(5)], os.path.join(DIAG,'dk_pipeline.png'), (15.5,2.4))

n={'a':hnode(0,'Client\nSubmits','terminal'),'b':hnode(1,'Review\nQueue'),'c':hnode(2,'Officer\nChecks Docs'),
   'd':hdiamond(3,'Compliance\nApproves?'),'e':hnode(4,'Verified\n(Unlocks)','terminal'),'r':hnode(3,'Rejected:\nResubmit',row=1)}
draw_flow(n, [('a','b'),('b','c'),('c','d'),('d','e','Yes'),('d','r','No')], os.path.join(DIAG,'dk_kyc.png'), (13.5,4.3))

n={'a':hnode(0,'Client\nRequests','terminal'),'b':hdiamond(1,'KYC + Funds\nValid?'),'c':hnode(2,'Pending\n(Held)'),
   'd':hdiamond(3,'Finance\nApproves?'),'e':hnode(4,'Paid /\nCompleted','terminal'),
   'b2':hnode(1,'Blocked:\nComplete KYC',row=1),'e2':hnode(3,'Rejected:\nHold Released',row=1)}
draw_flow(n, [('a','b'),('b','c','Yes'),('b','b2','No'),('c','d'),('d','e','Yes'),('d','e2','No')],
          os.path.join(DIAG,'dk_withdrawal.png'), (13.5,4.3))

n={'adm':(0,0,3.0,1.3,'ADMIN\nFull access — finance\napprovals, roles,\nsettings, all modules','group'),
   'agt':(3.8,0,3.0,1.3,'AGENT\nClients · Leads\nSupport · KYC review\n(no finance / roles)','group'),
   'aud':(1.9,-2.3,3.6,1.0,'AUDIT LOG\nEvery action recorded','terminal')}
draw_flow(n, [('adm','aud'),('agt','aud')], os.path.join(DIAG,'dk_roles.png'), (8.0,4.0))
print('Dark diagrams rendered.')

# ───────────────────────── deck ─────────────────────────
BG    = RGBColor(0x0A,0x0A,0x0A)
PANEL = RGBColor(0x14,0x17,0x1D)
ALT   = RGBColor(0x10,0x12,0x17)
cRED  = RGBColor(0xE1,0x1D,0x2E)
WHITE = RGBColor(0xFF,0xFF,0xFF)
LIGHT = RGBColor(0xCB,0xCF,0xD6)
MUTED = RGBColor(0x6B,0x70,0x7A)

prs = Presentation()
prs.slide_width = Inches(13.333); prs.slide_height = Inches(7.5)
W, Hh = prs.slide_width, prs.slide_height
BLANK = prs.slide_layouts[6]

def _noshadow(shp): shp.shadow.inherit = False

def base(title=None, kicker=None, footer=True):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, Hh)
    bg.fill.solid(); bg.fill.fore_color.rgb = BG; bg.line.fill.background(); _noshadow(bg)
    if title:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(0.52), Inches(0.13), Inches(0.62))
        bar.fill.solid(); bar.fill.fore_color.rgb = cRED; bar.line.fill.background(); _noshadow(bar)
        tb = s.shapes.add_textbox(Inches(0.9), Inches(0.4), Inches(11.9), Inches(1.0)); tf = tb.text_frame; tf.word_wrap = True
        if kicker:
            p = tf.paragraphs[0]; r = p.add_run(); r.text = kicker.upper()
            r.font.size = Pt(12); r.font.bold = True; r.font.color.rgb = cRED
            p2 = tf.add_paragraph()
        else:
            p2 = tf.paragraphs[0]
        r = p2.add_run(); r.text = title; r.font.size = Pt(30); r.font.bold = True; r.font.color.rgb = WHITE
    if footer:
        ft = s.shapes.add_textbox(Inches(0.6), Inches(7.04), Inches(12.1), Inches(0.32))
        p = ft.text_frame.paragraphs[0]
        r = p.add_run(); r.text = '27 Markets — CRM (Back-Office) Overview'; r.font.size = Pt(9); r.font.color.rgb = MUTED
        r2 = p.add_run(); r2.text = '     |     Prepared by Core Care Consultancy'; r2.font.size = Pt(9); r2.font.color.rgb = MUTED
    return s

def bullets(s, items, left=0.95, top=1.7, width=11.6, size=16, gap=12, headless=False):
    tb = s.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(5.0)); tf = tb.text_frame; tf.word_wrap = True
    for i, it in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph(); p.space_after = Pt(gap)
        rb = p.add_run(); rb.text = '▸ '; rb.font.color.rgb = cRED; rb.font.bold = True; rb.font.size = Pt(size)
        if headless:
            rbody = p.add_run(); rbody.text = it; rbody.font.color.rgb = LIGHT; rbody.font.size = Pt(size)
        else:
            head, body = it
            rh = p.add_run(); rh.text = head + ' — '; rh.font.bold = True; rh.font.color.rgb = WHITE; rh.font.size = Pt(size)
            rbody = p.add_run(); rbody.text = body; rbody.font.color.rgb = LIGHT; rbody.font.size = Pt(size)

def image(s, path, width_in, top_in, left_in=None):
    w = Inches(width_in)
    left = int((W - w) / 2) if left_in is None else Inches(left_in)
    s.shapes.add_picture(path, left, Inches(top_in), width=w)

def caption(s, text, top_in):
    tb = s.shapes.add_textbox(Inches(0.6), Inches(top_in), Inches(12.1), Inches(0.35))
    p = tb.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = text; r.font.italic = True; r.font.size = Pt(11); r.font.color.rgb = MUTED

def cell(c, text, fill, color, bold=False, size=11, align=PP_ALIGN.LEFT):
    c.fill.solid(); c.fill.fore_color.rgb = fill
    c.margin_left = Inches(0.1); c.margin_right = Inches(0.1); c.margin_top = Inches(0.03); c.margin_bottom = Inches(0.03)
    c.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf = c.text_frame; tf.word_wrap = True; p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run(); r.text = text; r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color

def table(s, headers, rows, left, top, width, col_widths, fontsize=11, rowh=0.4):
    gt = s.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(rowh * (len(rows) + 1))).table
    gt.first_row = False; gt.horz_banding = False
    for j, h in enumerate(headers):
        cell(gt.cell(0, j), h, cRED, WHITE, bold=True, size=fontsize + 1)
    for i, row in enumerate(rows):
        fill = PANEL if i % 2 == 0 else ALT
        for j, v in enumerate(row):
            cell(gt.cell(i + 1, j), str(v), fill, LIGHT, size=fontsize)
    for j, cw in enumerate(col_widths):
        gt.columns[j].width = Inches(cw)

# ── 1. TITLE ──
s = base(footer=False)
tb = s.shapes.add_textbox(Inches(0.9), Inches(2.0), Inches(11.5), Inches(3.4)); tf = tb.text_frame; tf.word_wrap = True
p = tf.paragraphs[0]; r = p.add_run(); r.text = '27 MARKETS'; r.font.size = Pt(46); r.font.bold = True; r.font.color.rgb = cRED
p = tf.add_paragraph(); r = p.add_run(); r.text = 'Trade Beyond Limits'; r.font.size = Pt(15); r.font.italic = True; r.font.color.rgb = MUTED
p = tf.add_paragraph(); p.space_before = Pt(18); r = p.add_run(); r.text = 'Client Relationship Management (CRM)'; r.font.size = Pt(28); r.font.bold = True; r.font.color.rgb = WHITE
p = tf.add_paragraph(); r = p.add_run(); r.text = 'Back-Office Platform — What It Includes & Why We Are Building It'; r.font.size = Pt(15); r.font.color.rgb = LIGHT
acc = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.95), Inches(5.55), Inches(2.3), Inches(0.06))
acc.fill.solid(); acc.fill.fore_color.rgb = cRED; acc.line.fill.background(); _noshadow(acc)
mb = s.shapes.add_textbox(Inches(0.9), Inches(5.75), Inches(11.5), Inches(1.2)); mf = mb.text_frame; mf.word_wrap = True
for txt in ['Prepared for:  27 Markets', 'Prepared by:  Core Care Consultancy', 'Date:  2026-06-22      |      Version 1.0']:
    p = mf.paragraphs[0] if txt.startswith('Prepared for') else mf.add_paragraph()
    r = p.add_run(); r.text = txt; r.font.size = Pt(12); r.font.color.rgb = LIGHT
if os.path.exists(LOGO):
    s.shapes.add_picture(LOGO, Inches(10.7), Inches(0.55), width=Inches(1.9))

# ── 2. WHAT IS THE CRM ──
s = base('What Is the CRM?', 'Overview')
tb = s.shapes.add_textbox(Inches(0.95), Inches(1.55), Inches(11.5), Inches(0.9)); tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = 'The private, secure back-office your team uses to run the business — the other half of the platform clients never see. One place to manage leads, clients, verification, money, support, and partners.'
r.font.size = Pt(15); r.font.color.rgb = LIGHT
image(s, os.path.join(DIAG, 'dk_modules.png'), width_in=7.1, top_in=2.55)
caption(s, 'The CRM modules, organised around the client.', 6.75)

# ── 3. WHY (part 1) ──
s = base('Why We Are Building It', 'The business case')
bullets(s, [
    ('One source of truth', 'Every lead, client, document, payment, and conversation in one place — no scattered spreadsheets.'),
    ('More conversions, more revenue', 'Leads are captured, assigned, and followed up systematically, so more visitors become funded clients.'),
    ('Regulatory compliance', 'A KYC review queue and a complete audit trail keep the business inspection-ready.'),
    ('Financial control & safety', 'No money moves without staff review — withdrawals are held and approved by finance.'),
], top=1.7, size=17, gap=16)

# ── 4. WHY (part 2) ──
s = base('Why We Are Building It', 'The business case (continued)')
bullets(s, [
    ('Operational visibility', 'Live dashboards show deposits, withdrawals, pending work, and team performance at a glance.'),
    ('Better client service', 'Staff see a client’s full history instantly, resolving requests faster and more personally.'),
    ('Accountability & security', 'Role-based access (Admin / Agent) plus an audit log — every action is recorded.'),
    ('Built to scale', 'From hundreds to thousands of clients on the same system; manual methods would break.'),
], top=1.7, size=17, gap=16)

# ── 5. WHAT IT INCLUDES ──
s = base('What the CRM Includes', 'Modules')
table(s, ['Module', 'What your team can do'], [
    ['Dashboard', 'Key numbers and pending work — new clients, KYC to review, withdrawals to approve, open tickets.'],
    ['Clients (360 View)', 'Full profile, accounts, balances, KYC status, transactions, notes, and activity history.'],
    ['Leads', 'A visual pipeline: capture, assign to agents, and track follow-ups to conversion.'],
    ['KYC / Compliance', 'Review identity documents in a queue and approve or reject them.'],
    ['Finance', 'See deposits and approve or reject withdrawals; controlled manual adjustments (Admin).'],
    ['Accounts', 'View trading accounts; suspend, activate, or change leverage.'],
    ['Support', 'A shared ticket inbox — assign, reply, and keep internal notes.'],
    ['Partners / IB', 'Partners and the clients they referred, with rebate tracking.'],
    ['Reports', 'Deposits, withdrawals, net flow, conversion, and agent performance.'],
    ['Staff & Settings', 'Manage staff, assign Admin/Agent access, review the audit log (Admin).'],
], left=0.95, top=1.6, width=11.45, col_widths=[2.6, 8.85], fontsize=11, rowh=0.5)

# ── 6. LEADS ──
s = base('Lead & Sales Pipeline', 'Growth')
tb = s.shapes.add_textbox(Inches(0.95), Inches(1.55), Inches(11.5), Inches(0.9)); tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = 'Leads are captured automatically from the website and demo requests, assigned to an agent, and tracked through clear stages until they fund and trade — so no opportunity is missed.'
r.font.size = Pt(15); r.font.color.rgb = LIGHT
image(s, os.path.join(DIAG, 'dk_pipeline.png'), width_in=12.0, top_in=3.4)
caption(s, 'From new lead to funded, active client.', 5.6)

# ── 7. CLIENT 360 ──
s = base('The Client 360° View', 'Everything on one screen')
bullets(s, [
    'Personal profile and contact details.',
    'All trading accounts with live balances.',
    'KYC verification status and uploaded documents.',
    'Full deposit, withdrawal, and transfer history.',
    'Internal staff notes and the complete activity timeline.',
], top=1.8, size=18, gap=16, headless=True)

# ── 8. KYC ──
s = base('KYC & Compliance Review', 'Safety & compliance')
image(s, os.path.join(DIAG, 'dk_kyc.png'), width_in=11.6, top_in=2.0)
caption(s, 'Submissions land in a queue; compliance approves or rejects. Withdrawals stay locked until verified.', 5.7)

# ── 9. FINANCE ──
s = base('Finance & Withdrawal Approvals', 'Money safety')
image(s, os.path.join(DIAG, 'dk_withdrawal.png'), width_in=11.6, top_in=2.0)
caption(s, 'Every withdrawal is checked, held, and approved by finance before any payout.', 5.7)

# ── 10. ROLES ──
s = base('Staff Roles & Access Control', 'Accountability')
image(s, os.path.join(DIAG, 'dk_roles.png'), width_in=7.0, top_in=1.7)
bullets(s, [
    'Admin: full access — finance approvals, staff/role management, settings, all modules.',
    'Agent: clients, leads, support, KYC review — no finance approvals or role changes.',
    'Audit log: every state-changing action recorded with who, what, and when.',
], left=0.95, top=5.05, width=11.6, size=13, gap=8, headless=True)

# ── 11. ROLLOUT ──
s = base('How We Will Build It', 'Phased delivery')
table(s, ['Phase', 'Delivered'], [
    ['1', 'Foundations — staff roles (Admin/Agent), secure access, and the admin shell.'],
    ['2', 'Dashboard with live numbers and recent activity.'],
    ['3', 'Clients (360 view + notes) and the KYC review queue.'],
    ['4', 'Finance (withdrawal approvals) and account administration.'],
    ['5', 'Leads pipeline and the support desk.'],
    ['6', 'Partners, reports, staff/settings, audit log, and final polish.'],
], left=0.95, top=1.8, width=11.45, col_widths=[1.3, 10.15], fontsize=13, rowh=0.62)

# ── 12. COSTS & GO-LIVE ──
s = base('Important Note — Costs & Go-Live', 'Please note')
bullets(s, [
    ('Current stage', 'This build runs on simulated data for demonstration. Live money movement is not active yet.'),
    ('Going live', 'Full launch once the trading platform, payment gateways, identity-verification provider, and licensing are connected.'),
    ('Third-party costs', 'Identity checks, email/SMS, and payment reconciliation rely on paid third-party services, billed by each provider.'),
    ('Future phases', 'Marketing automation, a full IB commission engine, and granular roles are planned later and quoted separately.'),
], top=1.7, size=16, gap=16)

# ── 13. SUMMARY ──
s = base('Summary', 'The control center')
tb = s.shapes.add_textbox(Inches(0.95), Inches(1.9), Inches(11.4), Inches(3.0)); tf = tb.text_frame; tf.word_wrap = True
r = tf.paragraphs[0].add_run()
r.text = ('The CRM is the control center of 27 Markets. It turns scattered, manual work into one secure system '
          'where leads convert faster, clients are served better, money moves only with proper approval, and the '
          'business stays compliant and fully auditable — with complete operational visibility today and a '
          'foundation that scales as 27 Markets grows.')
r.font.size = Pt(18); r.font.color.rgb = LIGHT
p = tf.add_paragraph(); p.space_before = Pt(26); r = p.add_run()
r.text = 'Prepared by Core Care Consultancy'; r.font.size = Pt(14); r.font.bold = True; r.font.color.rgb = cRED

out = os.path.join(BASE, '27-Markets-CRM-Presentation.pptx')
try:
    prs.save(out); print('Saved:', out)
except PermissionError:
    alt = os.path.join(BASE, '27-Markets-CRM-Presentation-v2.pptx'); prs.save(alt); print('Locked. Saved to:', alt)
